"""
Documents Router
Handles: PDF ↔ Word, PDF ↔ Excel, PDF ↔ PPT, PDF → Text/HTML,
         Word → PDF/Text/HTML, Excel → CSV/JSON, Markdown → HTML,
         HTML → PDF, CSV ↔ JSON ↔ XML
"""

import io
import csv
import json
import xml.dom.minidom as minidom
import xml.etree.ElementTree as ET
from pathlib import Path

from fastapi import APIRouter, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse

from utils.helpers import save_upload, output_path, cleanup, file_download_url

router = APIRouter()


# ═══════════════════════════════════════════════════════════════════════════════
#  PDF TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/pdf-to-text")
async def pdf_to_text(file: UploadFile = File(...)):
    """Extract plain text from PDF."""
    _check_ext(file.filename, [".pdf"])
    tmp = await save_upload(file)
    try:
        import pypdf
        reader = pypdf.PdfReader(str(tmp))
        text = "\n\n".join(
            page.extract_text() or "" for page in reader.pages
        ).strip()
        return {"text": text, "pages": len(reader.pages)}
    except ImportError:
        raise HTTPException(500, "pypdf not installed. Run: pip install pypdf")
    finally:
        cleanup(tmp)


@router.post("/pdf-to-html")
async def pdf_to_html(file: UploadFile = File(...)):
    """Extract text from PDF and wrap in HTML."""
    _check_ext(file.filename, [".pdf"])
    tmp = await save_upload(file)
    try:
        import pypdf
        reader = pypdf.PdfReader(str(tmp))
        pages_html = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            escaped = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            pages_html.append(f'<div class="page" id="page-{i}"><h2>Page {i}</h2><pre>{escaped}</pre></div>')
        html = f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="UTF-8"><title>PDF Export</title>
<style>
  body{{font-family:Georgia,serif;max-width:860px;margin:40px auto;padding:0 20px;background:#fafafa;color:#222}}
  .page{{background:#fff;padding:32px;margin:24px 0;border-radius:8px;box-shadow:0 2px 8px rgba(0,0,0,.08)}}
  h2{{color:#555;font-size:13px;text-transform:uppercase;letter-spacing:1px;margin-bottom:16px}}
  pre{{white-space:pre-wrap;font-family:inherit;line-height:1.7}}
</style></head>
<body>{''.join(pages_html)}</body></html>"""
        out = output_path(".html")
        out.write_text(html, encoding="utf-8")
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "pypdf not installed.")
    finally:
        cleanup(tmp)


@router.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """Convert PDF to DOCX using pdf2docx."""
    _check_ext(file.filename, [".pdf"])
    tmp = await save_upload(file)
    out = output_path(".docx")
    try:
        from pdf2docx import Converter
        cv = Converter(str(tmp))
        cv.convert(str(out))
        cv.close()
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "pdf2docx not installed. Run: pip install pdf2docx")
    finally:
        cleanup(tmp)


@router.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Extract tables from PDF into Excel using camelot or tabula."""
    _check_ext(file.filename, [".pdf"])
    tmp = await save_upload(file)
    out = output_path(".xlsx")
    try:
        import tabula
        import openpyxl
        dfs = tabula.read_pdf(str(tmp), pages="all", multiple_tables=True)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        for i, df in enumerate(dfs, 1):
            ws = wb.create_sheet(f"Table_{i}")
            ws.append(list(df.columns))
            for row in df.itertuples(index=False):
                ws.append(list(row))
        wb.save(str(out))
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "tabula-py / openpyxl not installed. Run: pip install tabula-py openpyxl")
    finally:
        cleanup(tmp)


@router.post("/pdf-to-images")
async def pdf_to_images(
    file: UploadFile = File(...),
    dpi: int = Form(150),
    fmt: str = Form("png")
):
    """Convert each PDF page to an image, return zip."""
    _check_ext(file.filename, [".pdf"])
    tmp = await save_upload(file)
    out = output_path(".zip")
    try:
        from pdf2image import convert_from_path
        import zipfile
        images = convert_from_path(str(tmp), dpi=dpi)
        with zipfile.ZipFile(str(out), "w") as zf:
            for i, img in enumerate(images, 1):
                buf = io.BytesIO()
                img.save(buf, format=fmt.upper())
                zf.writestr(f"page_{i:03d}.{fmt}", buf.getvalue())
        return {"url": file_download_url(out.name), "filename": out.name, "pages": len(images)}
    except ImportError:
        raise HTTPException(500, "pdf2image not installed. Run: pip install pdf2image")
    finally:
        cleanup(tmp)


# ═══════════════════════════════════════════════════════════════════════════════
#  WORD TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    """Convert DOCX to PDF using LibreOffice (soffice)."""
    _check_ext(file.filename, [".docx", ".doc"])
    tmp = await save_upload(file)
    out_dir = output_path("").parent
    try:
        import subprocess
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(tmp)],
            capture_output=True, timeout=60
        )
        # LibreOffice names output same stem + .pdf
        pdf_path = out_dir / (tmp.stem + ".pdf")
        if not pdf_path.exists():
            raise HTTPException(500, "Conversion failed. Ensure LibreOffice is installed.")
        # Rename to UUID
        final = output_path(".pdf")
        pdf_path.rename(final)
        return {"url": file_download_url(final.name), "filename": final.name}
    finally:
        cleanup(tmp)


@router.post("/word-to-text")
async def word_to_text(file: UploadFile = File(...)):
    """Extract plain text from DOCX."""
    _check_ext(file.filename, [".docx"])
    tmp = await save_upload(file)
    try:
        from docx import Document
        doc = Document(str(tmp))
        text = "\n".join(para.text for para in doc.paragraphs)
        return {"text": text}
    except ImportError:
        raise HTTPException(500, "python-docx not installed. Run: pip install python-docx")
    finally:
        cleanup(tmp)


@router.post("/word-to-html")
async def word_to_html(file: UploadFile = File(...)):
    """Convert DOCX to HTML using mammoth."""
    _check_ext(file.filename, [".docx"])
    tmp = await save_upload(file)
    out = output_path(".html")
    try:
        import mammoth
        with open(tmp, "rb") as f:
            result = mammoth.convert_to_html(f)
        html = f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>body{{font-family:Georgia,serif;max-width:860px;margin:40px auto;padding:0 20px;line-height:1.7}}</style>
</head><body>{result.value}</body></html>"""
        out.write_text(html, encoding="utf-8")
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "mammoth not installed. Run: pip install mammoth")
    finally:
        cleanup(tmp)


# ═══════════════════════════════════════════════════════════════════════════════
#  EXCEL TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/excel-to-csv")
async def excel_to_csv(
    file: UploadFile = File(...),
    sheet: int = Form(0)
):
    """Convert Excel sheet to CSV."""
    _check_ext(file.filename, [".xlsx", ".xls"])
    tmp = await save_upload(file)
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(tmp), data_only=True)
        ws = wb.worksheets[sheet]
        buf = io.StringIO()
        writer = csv.writer(buf)
        for row in ws.iter_rows(values_only=True):
            writer.writerow([str(c) if c is not None else "" for c in row])
        return {"csv": buf.getvalue(), "sheets": wb.sheetnames}
    except ImportError:
        raise HTTPException(500, "openpyxl not installed. Run: pip install openpyxl")
    finally:
        cleanup(tmp)


@router.post("/excel-to-json")
async def excel_to_json(
    file: UploadFile = File(...),
    sheet: int = Form(0)
):
    """Convert Excel sheet to JSON array using first row as headers."""
    _check_ext(file.filename, [".xlsx", ".xls"])
    tmp = await save_upload(file)
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(tmp), data_only=True)
        ws = wb.worksheets[sheet]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return {"data": [], "sheets": wb.sheetnames}
        headers = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(rows[0])]
        data = [dict(zip(headers, [str(c) if c is not None else "" for c in row])) for row in rows[1:]]
        return {"data": data, "sheets": wb.sheetnames}
    except ImportError:
        raise HTTPException(500, "openpyxl not installed.")
    finally:
        cleanup(tmp)


@router.post("/csv-to-excel")
async def csv_to_excel(file: UploadFile = File(...)):
    """Convert CSV to formatted Excel workbook."""
    _check_ext(file.filename, [".csv", ".txt"])
    tmp = await save_upload(file)
    out = output_path(".xlsx")
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        content = tmp.read_text(encoding="utf-8-sig")
        reader = csv.reader(io.StringIO(content))
        for i, row in enumerate(reader):
            ws.append(row)
            if i == 0:  # header row styling
                for cell in ws[1]:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill("solid", fgColor="2D6A4F")
                    cell.alignment = Alignment(horizontal="center")
        # Auto column widths
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=0)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)
        wb.save(str(out))
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "openpyxl not installed.")
    finally:
        cleanup(tmp)


@router.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    """Convert Excel to PDF via LibreOffice."""
    _check_ext(file.filename, [".xlsx", ".xls"])
    tmp = await save_upload(file)
    out_dir = output_path("").parent
    try:
        import subprocess
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(tmp)],
            capture_output=True, timeout=60
        )
        pdf_path = out_dir / (tmp.stem + ".pdf")
        if not pdf_path.exists():
            raise HTTPException(500, "LibreOffice conversion failed.")
        final = output_path(".pdf")
        pdf_path.rename(final)
        return {"url": file_download_url(final.name), "filename": final.name}
    finally:
        cleanup(tmp)


# ═══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT TOOLS
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/ppt-to-text")
async def ppt_to_text(file: UploadFile = File(...)):
    """Extract text from PowerPoint slides."""
    _check_ext(file.filename, [".pptx", ".ppt"])
    tmp = await save_upload(file)
    try:
        from pptx import Presentation
        prs = Presentation(str(tmp))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
            slides_text.append({"slide": i, "text": "\n".join(texts)})
        return {"slides": slides_text, "count": len(slides_text)}
    except ImportError:
        raise HTTPException(500, "python-pptx not installed. Run: pip install python-pptx")
    finally:
        cleanup(tmp)


@router.post("/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    """Convert PPTX to PDF via LibreOffice."""
    _check_ext(file.filename, [".pptx", ".ppt"])
    tmp = await save_upload(file)
    out_dir = output_path("").parent
    try:
        import subprocess
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf",
             "--outdir", str(out_dir), str(tmp)],
            capture_output=True, timeout=60
        )
        pdf_path = out_dir / (tmp.stem + ".pdf")
        if not pdf_path.exists():
            raise HTTPException(500, "LibreOffice not found. Install it for PPT→PDF.")
        final = output_path(".pdf")
        pdf_path.rename(final)
        return {"url": file_download_url(final.name), "filename": final.name}
    finally:
        cleanup(tmp)


# ═══════════════════════════════════════════════════════════════════════════════
#  CSV / JSON / XML CONVERSIONS
# ═══════════════════════════════════════════════════════════════════════════════

@router.post("/csv-to-json")
async def csv_to_json(file: UploadFile = File(...)):
    _check_ext(file.filename, [".csv", ".txt"])
    tmp = await save_upload(file)
    try:
        content = tmp.read_text(encoding="utf-8-sig")
        reader = csv.DictReader(io.StringIO(content))
        rows = list(reader)
        return {"data": rows, "count": len(rows)}
    finally:
        cleanup(tmp)


@router.post("/json-to-csv")
async def json_to_csv(file: UploadFile = File(...)):
    _check_ext(file.filename, [".json"])
    tmp = await save_upload(file)
    out = output_path(".csv")
    try:
        data = json.loads(tmp.read_text())
        if isinstance(data, dict):
            data = [data]
        if not data:
            raise HTTPException(400, "Empty JSON array")
        keys = list(data[0].keys())
        buf = io.StringIO()
        w = csv.DictWriter(buf, fieldnames=keys, extrasaction="ignore")
        w.writeheader()
        w.writerows(data)
        out.write_text(buf.getvalue(), encoding="utf-8")
        return {"url": file_download_url(out.name), "filename": out.name}
    finally:
        cleanup(tmp)


@router.post("/json-to-xml")
async def json_to_xml_file(file: UploadFile = File(...)):
    _check_ext(file.filename, [".json"])
    tmp = await save_upload(file)
    out = output_path(".xml")
    try:
        data = json.loads(tmp.read_text())
        xml_str = _dict_to_xml(data, "root")
        pretty = minidom.parseString(xml_str).toprettyxml(indent="  ")
        out.write_text(pretty, encoding="utf-8")
        return {"url": file_download_url(out.name), "filename": out.name}
    finally:
        cleanup(tmp)


@router.post("/xml-to-json")
async def xml_to_json_file(file: UploadFile = File(...)):
    _check_ext(file.filename, [".xml"])
    tmp = await save_upload(file)
    out = output_path(".json")
    try:
        tree = ET.parse(str(tmp))
        data = _xml_to_dict(tree.getroot())
        out.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        return {"url": file_download_url(out.name), "filename": out.name}
    finally:
        cleanup(tmp)


@router.post("/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    """Convert HTML file to PDF using WeasyPrint."""
    _check_ext(file.filename, [".html", ".htm"])
    tmp = await save_upload(file)
    out = output_path(".pdf")
    try:
        from weasyprint import HTML
        HTML(filename=str(tmp)).write_pdf(str(out))
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "weasyprint not installed. Run: pip install weasyprint")
    finally:
        cleanup(tmp)


@router.post("/markdown-to-html")
async def markdown_to_html(file: UploadFile = File(...)):
    """Convert Markdown file to styled HTML."""
    _check_ext(file.filename, [".md", ".markdown", ".txt"])
    tmp = await save_upload(file)
    out = output_path(".html")
    try:
        import markdown2
        md_text = tmp.read_text(encoding="utf-8")
        body = markdown2.markdown(md_text, extras=["fenced-code-blocks", "tables", "strike"])
        html = f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<style>
  body{{font-family:Georgia,serif;max-width:860px;margin:40px auto;padding:0 20px;line-height:1.75;color:#222}}
  pre{{background:#f4f4f4;padding:16px;border-radius:6px;overflow:auto}}
  code{{background:#f0f0f0;padding:2px 5px;border-radius:3px}}
  blockquote{{border-left:4px solid #ccc;margin:0;padding-left:16px;color:#666}}
  table{{border-collapse:collapse;width:100%}}
  th,td{{border:1px solid #ddd;padding:8px 12px}}
  th{{background:#f4f4f4}}
</style></head><body>{body}</body></html>"""
        out.write_text(html, encoding="utf-8")
        return {"url": file_download_url(out.name), "filename": out.name}
    except ImportError:
        raise HTTPException(500, "markdown2 not installed. Run: pip install markdown2")
    finally:
        cleanup(tmp)


# ═══════════════════════════════════════════════════════════════════════════════
#  PRIVATE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _check_ext(filename: str, allowed: list):
    if not filename:
        raise HTTPException(400, "No filename provided")
    ext = Path(filename).suffix.lower()
    if ext not in allowed:
        raise HTTPException(400, f"Expected {allowed}, got '{ext}'")


def _dict_to_xml(data, tag="item") -> str:
    safe = lambda s: str(s).replace(" ", "_").replace("/", "_") or "item"
    if isinstance(data, dict):
        children = "".join(_dict_to_xml(v, safe(k)) for k, v in data.items())
        return f"<{tag}>{children}</{tag}>"
    elif isinstance(data, list):
        return "".join(_dict_to_xml(item, "item") for item in data)
    else:
        val = str(data).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
        return f"<{tag}>{val}</{tag}>"


def _xml_to_dict(element):
    result = {}
    for child in element:
        child_data = _xml_to_dict(child)
        if child.tag in result:
            if not isinstance(result[child.tag], list):
                result[child.tag] = [result[child.tag]]
            result[child.tag].append(child_data)
        else:
            result[child.tag] = child_data
    return result if result else (element.text or "")
